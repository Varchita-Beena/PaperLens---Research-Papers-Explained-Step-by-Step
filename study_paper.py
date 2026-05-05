import json
import os
import re
import sys

from baseline_summary import call_openai
from baseline_summary import extract_text_from_pdf
from download_paper import step1_download_paper


def save_json(path, data):
    folder = os.path.dirname(path)
    if folder:
        os.makedirs(folder, exist_ok=True)

    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)


def clean_json(text):
    text = text.strip()
    text = text.replace("```json", "")
    text = text.replace("```", "")
    text = text.strip()

    try:
        return json.loads(text)
    except Exception:
        return {"raw_output": text}


def find_section(text, possible_names):
    lines = text.splitlines()
    start = -1
    end = len(lines)

    for i, line in enumerate(lines):
        clean_line = clean_heading(line)

        for name in possible_names:
            if clean_line == name.lower() or clean_line.startswith(name.lower()):
                start = i
                break

        if start != -1:
            break

    if start == -1:
        return ""

    section_names = [
        "abstract",
        "introduction",
        "related work",
        "background",
        "method",
        "methodology",
        "approach",
        "model",
        "experiments",
        "experimental setup",
        "datasets",
        "results",
        "discussion",
        "conclusion",
        "references",
    ]

    for i in range(start + 1, len(lines)):
        clean_line = clean_heading(lines[i])
        for name in section_names:
            if clean_line == name or clean_line.startswith(name):
                end = i
                return "\n".join(lines[start:end])

    return "\n".join(lines[start:end])


def clean_heading(line):
    line = line.strip().lower()
    line = re.sub(r"^\d+(\.\d+)*\s+", "", line)
    return line


def make_sections(paper_text):
    sections = {}

    sections["abstract"] = find_section(paper_text, ["abstract"])
    sections["introduction"] = find_section(paper_text, ["introduction", "1 introduction"])
    sections["related_work"] = find_section(
        paper_text,
        ["related work", "background", "2 related work", "2 background"],
    )
    sections["methodology"] = find_section(
        paper_text,
        ["method", "methodology", "approach", "model", "proposed method"],
    )
    sections["datasets"] = find_section(paper_text, ["datasets", "data", "dataset"])
    sections["experiments"] = find_section(
        paper_text,
        ["experiments", "experimental setup", "evaluation"],
    )
    sections["results"] = find_section(paper_text, ["results", "discussion"])
    sections["conclusion"] = find_section(paper_text, ["conclusion", "conclusions"])

    if sections["abstract"] == "":
        sections["abstract"] = paper_text[:4000]

    if sections["introduction"] == "":
        sections["introduction"] = paper_text[:10000]

    return sections


def short(text, characters=18000):
    return text[:characters]


def find_text_with_keywords(paper_text, keywords):
    chunks = paper_text.split("\n\n")
    selected = []

    for chunk in chunks:
        lower_chunk = chunk.lower()

        for keyword in keywords:
            if keyword.lower() in lower_chunk:
                selected.append(chunk)
                break

        if len("\n\n".join(selected)) > 18000:
            break

    return "\n\n".join(selected)


def ask_step(step_name, prompt, api_key, state):
    print("Running", step_name)

    llm_output = call_openai(prompt, api_key)
    parsed_output = clean_json(llm_output)

    state[step_name] = parsed_output
    state["raw_" + step_name] = llm_output

    save_json("PaperLens/outputs/" + step_name + ".json", parsed_output)
    save_json("PaperLens/outputs/current_state.json", state)

    return parsed_output


def read_file(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def make_ppt_prompt(state):
    skill_text = read_file("PaperLens/skills/ppt_writer/skill.md")

    return skill_text + """

Final paper study JSON:
""" + json.dumps(state["step10_final_combined"])


def save_ppt_text(slides, path):
    lines = []

    for i, slide in enumerate(slides, start=1):
        lines.append("Slide " + str(i) + ": " + slide.get("title", ""))

        for bullet in slide.get("bullets", []):
            lines.append("- " + str(bullet))

        lines.append("Speaker notes: " + slide.get("speaker_notes", ""))
        lines.append("Visual suggestion: " + slide.get("visual_suggestion", ""))
        lines.append("")

    folder = os.path.dirname(path)
    if folder:
        os.makedirs(folder, exist_ok=True)

    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))


def make_pptx(slides, output_path):
    try:
        from pptx import Presentation
    except Exception:
        print("PPTX not created because python-pptx is not installed.")
        print("This script is using Python:", sys.executable)
        print("Install it with:")
        print(sys.executable + " -m pip install python-pptx")
        return

    folder = os.path.dirname(output_path)
    if folder:
        os.makedirs(folder, exist_ok=True)

    presentation = Presentation()

    for slide_data in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = str(slide_data.get("title", ""))

        bullets = slide_data.get("bullets", [])
        body = slide.placeholders[1]
        text_frame = body.text_frame
        text_frame.clear()

        for i, bullet in enumerate(bullets[:6]):
            if i == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            paragraph.text = str(bullet)
            paragraph.level = 0

        notes = slide.notes_slide.notes_text_frame
        notes.text = str(slide_data.get("speaker_notes", ""))

    presentation.save(output_path)


def step3_problem_prompt(state):
    return """
You are explaining a research paper to students in very simple language.

Task:
Explain the problem statement of the paper.

Important:
Sometimes papers use a big technical sentence like "extracting protein-protein interactions".
Do not just repeat that sentence. Explain what it actually means.
Use simple language. Give analogies if useful.

Explain:
1. What is the paper about?
2. What exact problem is the paper trying to solve?
3. What does that problem statement mean in simple words?
4. Why is this problem significant?
5. What should a reader expect to learn in this paper?

Return only JSON:
{
  "what_is_the_paper_about": "",
  "problem_statement": "",
  "simple_meaning_of_problem": "",
  "significance": "",
  "what_to_expect_in_paper": "",
  "student_friendly_analogy": ""
}

Paper abstract:
""" + short(state["sections"]["abstract"]) + """

Paper introduction:
""" + short(state["sections"]["introduction"])


def step4_intro_prompt(state):
    return """
Explain the introduction section of this paper in simple language.

Use the previous problem explanation while explaining the introduction.

Explain:
1. What background does the introduction give?
2. What motivation does it give?
3. What gap or limitation does it point out?
4. What contributions does the paper claim?
5. How does the introduction prepare the reader for the rest of the paper?

Return only JSON:
{
  "background": "",
  "motivation": "",
  "gap_or_limitation": "",
  "claimed_contributions": [],
  "how_intro_sets_up_paper": "",
  "simple_explanation": ""
}

Previous output:
""" + json.dumps(state["step3_problem"]) + """

Introduction section:
""" + short(state["sections"]["introduction"])


def step5_related_work_prompt(state):
    return """
Explain the related work or background section in simple language.

Use the previous outputs to explain how earlier work connects to the paper's problem.

Explain:
1. What approaches existed before this paper?
2. What were their strengths?
3. What were their weaknesses or gaps?
4. How does this paper position itself against previous work?
5. Why was a new approach needed?

Return only JSON:
{
  "previous_approaches": [],
  "strengths_of_previous_work": [],
  "weaknesses_or_gaps": [],
  "how_this_paper_is_different": "",
  "why_new_approach_needed": "",
  "simple_explanation": ""
}

Previous problem output:
""" + json.dumps(state["step3_problem"]) + """

Previous introduction output:
""" + json.dumps(state["step4_introduction"]) + """

Related work section:
""" + short(state["sections"]["related_work"])


def step6_methodology_prompt(state):
    return """
Explain the methodology of the paper in very simple language.

If the methodology has many parts, explain each part separately.
If there are equations, explain what they mean conceptually instead of only repeating symbols.
If there is an architecture, explain the flow of information through it.

Explain:
1. What is the proposed method?
2. What are the main components?
3. How does each component work?
4. How do the components connect?
5. What is new or important in the method?
6. What are the difficult technical parts, simplified?

Return only JSON:
{
  "method_overview": "",
  "method_parts": [],
  "technical_terms_explained": [],
  "equations_explained": [],
  "architecture_explained": "",
  "what_is_new": "",
  "simple_explanation": ""
}

Previous outputs:
""" + json.dumps({
        "problem": state["step3_problem"],
        "introduction": state["step4_introduction"],
        "related_work": state["step5_related_work"],
    }) + """

Methodology section:
""" + short(state["sections"]["methodology"])


def step7_datasets_prompt(state):
    return """
Explain the datasets used in this paper in simple language.

Explain:
1. Which datasets are used?
2. What does each dataset contain?
3. Why are these datasets suitable for this problem?
4. How are datasets split or prepared, if mentioned?
5. Any limitations of the datasets.

Return only JSON:
{
  "datasets": [],
  "what_each_dataset_contains": [],
  "why_datasets_are_used": "",
  "data_preparation": "",
  "dataset_limitations": [],
  "simple_explanation": ""
}

Previous method output:
""" + json.dumps(state["step6_methodology"]) + """

Dataset and experiment text:
""" + short(state["dataset_text"])


def step8_experiments_prompt(state):
    return """
Explain the experiments designed and performed in the paper.

Explain the whole experimental setting:
1. What was being tested?
2. What baselines were compared?
3. What metrics were used?
4. What experimental setup was used?
5. What ablations or variations were tested, if any?
6. Why these experiments make sense for the problem.

Return only JSON:
{
  "what_was_tested": "",
  "baselines": [],
  "metrics": [],
  "experimental_setup": "",
  "ablations_or_variations": [],
  "why_experiments_make_sense": "",
  "simple_explanation": ""
}

Previous outputs:
""" + json.dumps({
        "problem": state["step3_problem"],
        "methodology": state["step6_methodology"],
        "datasets": state["step7_datasets"],
    }) + """

Experiment section:
""" + short(state["sections"]["experiments"])


def step9_results_prompt(state):
    return """
Explain the results of the paper in simple language.

If the results section is large, separate the explanation into parts.

Explain:
1. What were the main results?
2. Which method performed best?
3. What do the numbers mean in simple words?
4. What results support the paper's main claim?
5. Were there surprising or weak results?
6. What can we conclude from the results?

Return only JSON:
{
  "main_results": [],
  "best_performing_method": "",
  "numbers_explained_simply": "",
  "evidence_for_main_claim": "",
  "surprising_or_weak_results": [],
  "result_parts": [],
  "simple_conclusion_from_results": ""
}

Previous outputs:
""" + json.dumps({
        "problem": state["step3_problem"],
        "methodology": state["step6_methodology"],
        "datasets": state["step7_datasets"],
        "experiments": state["step8_experiments"],
    }) + """

Results section:
""" + short(state["sections"]["results"] + "\n" + state["sections"]["conclusion"])


def step10_final_prompt(state):
    return """
Combine all previous step outputs into one final structured JSON.

This is not a new summary from scratch.
Use the previous outputs as the main source.

Return only JSON:
{
  "paper_study": {
    "problem_statement": {},
    "introduction": {},
    "related_work": {},
    "methodology": {},
    "datasets": {},
    "experiments": {},
    "results": {}
  },
  "overall_simple_explanation": "",
  "key_takeaways_for_students": [],
  "possible_ppt_slides": [
    {
      "slide_title": "",
      "bullet_points": []
    }
  ]
}

Previous outputs:
""" + json.dumps({
        "problem_statement": state["step3_problem"],
        "introduction": state["step4_introduction"],
        "related_work": state["step5_related_work"],
        "methodology": state["step6_methodology"],
        "datasets": state["step7_datasets"],
        "experiments": state["step8_experiments"],
        "results": state["step9_results"],
    })


def run_full_study(paper_input, api_key):
    state = {}
    state["paper_input"] = paper_input

    print("Running step1_download_paper")
    step1 = step1_download_paper(paper_input)
    state["step1_download"] = step1
    state["pdf_path"] = step1["pdf_path"]

    print("Running step2_extract_sections")
    paper_text = extract_text_from_pdf(state["pdf_path"])
    state["paper_text_characters"] = len(paper_text)
    state["sections"] = make_sections(paper_text)

    dataset_text = state["sections"]["datasets"]
    dataset_text = dataset_text + "\n" + state["sections"]["experiments"]
    dataset_text = dataset_text + "\n" + state["sections"]["results"]

    if len(dataset_text.strip()) < 1000:
        dataset_text = find_text_with_keywords(
            paper_text,
            [
                "dataset",
                "data set",
                "training data",
                "test set",
                "benchmark",
                "corpus",
                "translation task",
            ],
        )

    state["dataset_text"] = dataset_text

    save_json("PaperLens/outputs/sections.json", state["sections"])

    ask_step("step3_problem", step3_problem_prompt(state), api_key, state)
    ask_step("step4_introduction", step4_intro_prompt(state), api_key, state)
    ask_step("step5_related_work", step5_related_work_prompt(state), api_key, state)
    ask_step("step6_methodology", step6_methodology_prompt(state), api_key, state)
    ask_step("step7_datasets", step7_datasets_prompt(state), api_key, state)
    ask_step("step8_experiments", step8_experiments_prompt(state), api_key, state)
    ask_step("step9_results", step9_results_prompt(state), api_key, state)
    ask_step("step10_final_combined", step10_final_prompt(state), api_key, state)

    ppt_output = ask_step("step11_ppt_slides", make_ppt_prompt(state), api_key, state)
    slides = ppt_output.get("slides", [])
    save_ppt_text(slides, "PaperLens/outputs/presentation_text.txt")
    make_pptx(slides, "PaperLens/outputs/presentation.pptx")

    save_json("PaperLens/outputs/full_state.json", state)

    return state


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python3 PaperLens/study_paper.py paper name or paper link OPENAI_API_KEY")
        sys.exit(1)

    paper_input = " ".join(sys.argv[1:-1])
    api_key = sys.argv[-1]

    run_full_study(paper_input, api_key)

    print("Full paper study complete")
    print("All outputs saved in: PaperLens/outputs")
    print("Final JSON: PaperLens/outputs/step10_final_combined.json")
    print("PPT text: PaperLens/outputs/presentation_text.txt")
    print("PPTX: PaperLens/outputs/presentation.pptx")
    print("Full state: PaperLens/outputs/full_state.json")
